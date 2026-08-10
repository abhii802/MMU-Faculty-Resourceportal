#!/usr/bin/env python3
"""
CS Topic Presentation Generator
================================
Generates a modern, light-themed PPTX presentation for any Computer Science topic.
Includes diagrams, real-world examples, and beginner-friendly explanations.

Requirements:
    pip install python-pptx openai

Usage:
    python cs_pptx_generator.py "Data Structures"

    # With custom output path:
    python cs_pptx_generator.py "Machine Learning" --output ml_presentation.pptx

    # Without API key (uses built-in templates):
    python cs_pptx_generator.py "Arrays" --no-api

Environment:
    OPENAI_API_KEY - Required for AI-generated content (recommended)
"""

import os
import sys
import argparse
import json
from typing import List, Dict
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# =============================================================================
# CONFIGURATION - Modern Light Theme
# =============================================================================

@dataclass
class Theme:
    """Modern Light Theme Colors"""
    bg_white: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)
    bg_light: RGBColor = RGBColor(0xF8, 0xFA, 0xFC)
    bg_card: RGBColor = RGBColor(0xF0, 0xF4, 0xF8)
    primary: RGBColor = RGBColor(0x3B, 0x82, 0xF6)
    primary_dark: RGBColor = RGBColor(0x1D, 0x4E, 0xD8)
    primary_light: RGBColor = RGBColor(0xDB, 0xEA, 0xFE)
    accent_teal: RGBColor = RGBColor(0x14, 0xB8, 0xA6)
    accent_green: RGBColor = RGBColor(0x10, 0xB9, 0x81)
    accent_orange: RGBColor = RGBColor(0xF9, 0x73, 0x16)
    accent_purple: RGBColor = RGBColor(0x8B, 0x5C, 0xF6)
    accent_red: RGBColor = RGBColor(0xEF, 0x44, 0x44)
    text_dark: RGBColor = RGBColor(0x1E, 0x29, 0x3B)
    text_medium: RGBColor = RGBColor(0x47, 0x55, 0x69)
    text_light: RGBColor = RGBColor(0x94, 0xA3, 0xB8)
    text_white: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)

THEME = Theme()

# =============================================================================
# CONTENT GENERATION
# =============================================================================

class ContentGenerator:
    """Generates structured content for CS presentations."""

    def __init__(self, use_api: bool = True):
        self.use_api = use_api and bool(os.getenv("OPENAI_API_KEY"))
        self.client = None
        if self.use_api:
            try:
                import openai
                self.client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
            except ImportError:
                print("Warning: openai not installed. Using template mode.")
                self.use_api = False

    def generate(self, topic: str) -> Dict:
        if self.use_api:
            return self._generate_ai(topic)
        return self._generate_template(topic)

    def _generate_ai(self, topic: str) -> Dict:
        prompt = f"""Create a beginner-friendly Computer Science presentation about "{topic}".

Return ONLY a JSON object with this exact structure:
{{
    "title": "Main title (catchy but clear)",
    "subtitle": "One-line description",
    "introduction": {{
        "heading": "What is {topic}?",
        "points": ["3-4 simple sentences explaining the concept"]
    }},
    "why_it_matters": {{
        "heading": "Why Should You Care?",
        "points": ["3 real-world applications with simple explanations"]
    }},
    "core_concepts": {{
        "heading": "Core Concepts",
        "concepts": [
            {{"name": "Concept Name", "description": "Simple 1-2 line explanation", "example": "Real world analogy"}},
            {{"name": "Concept Name", "description": "Simple 1-2 line explanation", "example": "Real world analogy"}},
            {{"name": "Concept Name", "description": "Simple 1-2 line explanation", "example": "Real world analogy"}}
        ]
    }},
    "how_it_works": {{
        "heading": "How Does It Work?",
        "steps": ["5-6 simple steps explaining the process"]
    }},
    "real_world_examples": {{
        "heading": "Real-World Examples",
        "examples": [
            {{"domain": "Example Domain", "description": "How {topic} is used here", "benefit": "What problem it solves"}},
            {{"domain": "Example Domain", "description": "How {topic} is used here", "benefit": "What problem it solves"}},
            {{"domain": "Example Domain", "description": "How {topic} is used here", "benefit": "What problem it solves"}}
        ]
    }},
    "comparison": {{
        "heading": "Pros & Cons",
        "pros": ["3 advantages"],
        "cons": ["3 limitations or challenges"]
    }},
    "summary": {{
        "heading": "Key Takeaways",
        "points": ["5 bullet points summarizing the most important ideas"]
    }}
}}

Rules:
- Use extremely simple language suitable for beginners
- Every concept MUST have a real-world analogy or example
- Keep descriptions under 25 words each
- Focus on practical understanding over theory
- Use analogies from everyday life (cooking, organizing, traffic, etc.)
"""
        try:
            response = self.client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are an expert CS educator who explains complex topics simply."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=2500
            )
            content = response.choices[0].message.content
            if "```json" in content:
                content = content.split("```json")[1].split("```")[0]
            elif "```" in content:
                content = content.split("```")[1].split("```")[0]
            return json.loads(content.strip())
        except Exception as e:
            print(f"AI generation failed: {e}. Using template mode.")
            return self._generate_template(topic)

    def _generate_template(self, topic: str) -> Dict:
        return {
            "title": f"Understanding {topic}",
            "subtitle": f"A Beginner's Guide to {topic} in Computer Science",
            "introduction": {
                "heading": f"What is {topic}?",
                "points": [
                    f"{topic} is a fundamental concept in Computer Science that helps solve real problems.",
                    "Think of it like a toolbox that programmers use to build better software.",
                    "Understanding this topic will make you a better problem solver.",
                    "It is used by companies like Google, Amazon, and Netflix every day."
                ]
            },
            "why_it_matters": {
                "heading": "Why Should You Care?",
                "points": [
                    f"Every app on your phone uses {topic} in some way to work smoothly.",
                    f"Learning {topic} helps you write faster and more efficient code.",
                    f"It is asked in almost every tech interview at big companies."
                ]
            },
            "core_concepts": {
                "heading": "Core Concepts",
                "concepts": [
                    {"name": "Basic Building Block", "description": "The simplest unit of this topic that everything else is built from.", "example": "Like bricks in a house - small but essential."},
                    {"name": "Operations", "description": "The actions we can perform using this topic.", "example": "Like adding, removing, or finding items in a list."},
                    {"name": "Use Cases", "description": "Situations where this topic shines and solves problems best.", "example": "Like using a calculator for math instead of counting on fingers."}
                ]
            },
            "how_it_works": {
                "heading": "How Does It Work?",
                "steps": [
                    "Identify the problem you need to solve",
                    "Choose the right approach based on your needs",
                    "Apply the basic operations step by step",
                    "Check if the result solves your problem",
                    "Optimize if needed for better performance"
                ]
            },
            "real_world_examples": {
                "heading": "Real-World Examples",
                "examples": [
                    {"domain": "Social Media", "description": f"Platforms like Instagram use {topic} to organize your feed.", "benefit": "Shows you relevant content quickly without waiting."},
                    {"domain": "Online Shopping", "description": f"Amazon uses {topic} to manage millions of products.", "benefit": "Helps you find and buy items in seconds."},
                    {"domain": "Navigation Apps", "description": f"Google Maps uses {topic} to find the fastest route.", "benefit": "Saves you time and fuel every day."}
                ]
            },
            "comparison": {
                "heading": "Pros & Cons",
                "pros": ["Easy to learn and understand", "Widely used in industry", "Makes code more efficient"],
                "cons": ["Can be tricky at first", "Requires practice to master", "Not perfect for every situation"]
            },
            "summary": {
                "heading": "Key Takeaways",
                "points": [
                    f"{topic} is essential for writing good software.",
                    "Start with the basics and build your understanding gradually.",
                    "Always think about real-world applications while learning.",
                    "Practice is the key to mastering any CS topic.",
                    "This concept is used by top tech companies worldwide."
                ]
            }
        }


# =============================================================================
# PRESENTATION BUILDER
# =============================================================================

class PresentationBuilder:
    """Builds a modern light-themed PPTX presentation."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _add_background_shape(self, slide, color: RGBColor, left=0, top=0, width=None, height=None):
        if width is None:
            width = self.prs.slide_width
        if height is None:
            height = self.prs.slide_height
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        spTree = slide.shapes._spTree
        sp = shape._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        return shape

    def _add_text_box(self, slide, left, top, width, height, text, 
                      font_size: int = 18, bold: bool = False, 
                      color: RGBColor = None, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
        if color:
            p.font.color.rgb = color
        else:
            p.font.color.rgb = THEME.text_dark
        return txBox

    def _add_diagram_shape(self, slide, left, top, width, height, 
                          text: str, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                          fill_color: RGBColor = None, text_color: RGBColor = None,
                          font_size: int = 14):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color or THEME.primary
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = text_color or THEME.text_white
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        return shape

    def create_title_slide(self, title: str, subtitle: str):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(0.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = THEME.primary
        top_bar.line.fill.background()

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = THEME.primary_light
        circle.line.fill.background()
        spTree = slide.shapes._spTree
        sp = circle._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4), Inches(4), Inches(4))
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = THEME.primary_light
        circle2.line.fill.background()
        spTree = slide.shapes._spTree
        sp = circle2._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        self._add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                          title, font_size=54, bold=True, color=THEME.primary_dark, align=PP_ALIGN.CENTER)

        self._add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
                          subtitle, font_size=24, color=THEME.text_medium, align=PP_ALIGN.CENTER)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     self.prs.slide_width / 2 - Inches(1.5), Inches(4.5), 
                                     Inches(3), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = THEME.primary
        line.line.fill.background()

        self._add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
                          "Computer Science Fundamentals", font_size=12, 
                          color=THEME.text_light, align=PP_ALIGN.CENTER)

    def create_content_slide(self, heading: str, points: List[str], accent_color: RGBColor = None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        accent = accent_color or THEME.primary

        self._add_background_shape(slide, THEME.bg_white)

        left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), self.prs.slide_height)
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = accent
        left_bar.line.fill.background()

        self._add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8),
                          heading, font_size=36, bold=True, color=THEME.primary_dark)

        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), 
                                            Inches(1.5), Inches(0.06))
        underline.fill.solid()
        underline.fill.fore_color.rgb = accent
        underline.line.fill.background()

        y_pos = Inches(1.8)
        for point in points:
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y_pos + Inches(0.08), 
                                           Inches(0.12), Inches(0.12))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = accent
            bullet.line.fill.background()

            self._add_text_box(slide, Inches(1.3), y_pos, Inches(10.5), Inches(0.8),
                              point, font_size=18, color=THEME.text_medium)
            y_pos += Inches(0.9)

    def create_concept_slide(self, heading: str, concepts: List[Dict]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_background_shape(slide, THEME.bg_light)

        self._add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
                          heading, font_size=34, bold=True, color=THEME.primary_dark)

        colors = [THEME.accent_teal, THEME.accent_green, THEME.accent_orange]
        card_width = Inches(3.8)
        start_x = Inches(0.6)

        for i, concept in enumerate(concepts[:3]):
            x = start_x + i * (card_width + Inches(0.3))

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), 
                                         card_width, Inches(2.2))
            card.fill.solid()
            card.fill.fore_color.rgb = THEME.bg_white
            card.line.color.rgb = THEME.bg_card
            card.line.width = Pt(1)

            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.3), 
                                            card_width, Inches(0.08))
            accent.fill.solid()
            accent.fill.fore_color.rgb = colors[i]
            accent.line.fill.background()

            self._add_text_box(slide, x + Inches(0.15), Inches(1.45), card_width - Inches(0.3), Inches(0.4),
                              concept["name"], font_size=18, bold=True, color=colors[i])

            self._add_text_box(slide, x + Inches(0.15), Inches(1.85), card_width - Inches(0.3), Inches(0.7),
                              concept["description"], font_size=13, color=THEME.text_medium)

            example_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                                x + Inches(0.15), Inches(2.6), 
                                                card_width - Inches(0.3), Inches(0.7))
            example_box.fill.solid()
            example_box.fill.fore_color.rgb = THEME.primary_light
            example_box.line.fill.background()

            self._add_text_box(slide, x + Inches(0.25), Inches(2.65), 
                              card_width - Inches(0.5), Inches(0.6),
                              f"Example: {concept['example']}", font_size=12, 
                              color=THEME.primary_dark)

        self._draw_concept_map(slide, concepts, start_y=Inches(4.0))

    def _draw_concept_map(self, slide, concepts, start_y):
        center_x = self.prs.slide_width / 2

        center_shape = self._add_diagram_shape(slide, center_x - Inches(1), start_y, 
                                              Inches(2), Inches(0.6), 
                                              "Core Topic", MSO_SHAPE.ROUNDED_RECTANGLE,
                                              THEME.primary, THEME.text_white, 14)

        colors = [THEME.accent_teal, THEME.accent_green, THEME.accent_orange]
        positions = [Inches(2), self.prs.slide_width / 2 - Inches(1), Inches(9)]

        for i, concept in enumerate(concepts[:3]):
            x_pos = positions[i]
            y_pos = start_y + Inches(1.2)

            node = self._add_diagram_shape(slide, x_pos, y_pos, Inches(2.2), Inches(0.55),
                                          concept["name"], MSO_SHAPE.ROUNDED_RECTANGLE,
                                          colors[i], THEME.text_white, 12)

            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                             center_shape.left + center_shape.width / 2,
                                             center_shape.top + center_shape.height,
                                             node.left + node.width / 2,
                                             node.top)
            line.line.color.rgb = THEME.text_light
            line.line.width = Pt(2)

    def create_flowchart_slide(self, heading: str, steps: List[str]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_background_shape(slide, THEME.bg_white)

        self._add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
                          heading, font_size=34, bold=True, color=THEME.primary_dark)

        box_width = Inches(3.5)
        box_height = Inches(0.6)
        start_x = self.prs.slide_width / 2 - box_width / 2
        start_y = Inches(1.4)
        gap = Inches(0.85)

        colors = [THEME.primary, THEME.accent_teal, THEME.accent_green, 
                 THEME.accent_orange, THEME.accent_purple, THEME.primary_dark]

        for i, step in enumerate(steps[:6]):
            y = start_y + i * gap

            if i % 2 == 0:
                x = start_x - Inches(0.5)
            else:
                x = start_x + Inches(0.5)

            num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                               self.prs.slide_width / 2 - Inches(0.2), 
                                               y + Inches(0.05), Inches(0.4), Inches(0.4))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = colors[i % len(colors)]
            num_circle.line.fill.background()

            num_tf = num_circle.text_frame
            num_tf.word_wrap = False
            num_p = num_tf.paragraphs[0]
            num_p.text = str(i + 1)
            num_p.font.size = Pt(16)
            num_p.font.bold = True
            num_p.font.color.rgb = THEME.text_white
            num_p.alignment = PP_ALIGN.CENTER

            box = self._add_diagram_shape(slide, x, y, box_width, box_height,
                                           step, MSO_SHAPE.ROUNDED_RECTANGLE,
                                           colors[i % len(colors)], THEME.text_white, 12)

            if i < len(steps[:6]) - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                              self.prs.slide_width / 2 - Inches(0.15),
                                              y + box_height + Inches(0.05),
                                              Inches(0.3), Inches(0.2))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = THEME.text_light
                arrow.line.fill.background()

    def create_examples_slide(self, heading: str, examples: List[Dict]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_background_shape(slide, THEME.bg_light)

        self._add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
                          heading, font_size=34, bold=True, color=THEME.primary_dark)

        card_width = Inches(3.9)
        card_height = Inches(2.8)
        start_x = Inches(0.5)
        start_y = Inches(1.3)

        icons = ["WEB", "SHOP", "MAP"]
        colors = [THEME.accent_teal, THEME.accent_orange, THEME.accent_purple]

        for i, example in enumerate(examples[:3]):
            x = start_x + i * (card_width + Inches(0.3))

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, 
                                         card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = THEME.bg_white
            card.line.color.rgb = THEME.bg_card

            icon_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                                  x + card_width / 2 - Inches(0.3), 
                                                  start_y + Inches(0.2), 
                                                  Inches(0.6), Inches(0.6))
            icon_circle.fill.solid()
            icon_circle.fill.fore_color.rgb = colors[i]
            icon_circle.line.fill.background()

            self._add_text_box(slide, x + card_width / 2 - Inches(0.3), 
                              start_y + Inches(0.3), Inches(0.6), Inches(0.4),
                              icons[i], font_size=14, bold=True, 
                              color=THEME.text_white, align=PP_ALIGN.CENTER)

            self._add_text_box(slide, x + Inches(0.15), start_y + Inches(0.9), 
                              card_width - Inches(0.3), Inches(0.4),
                              example["domain"], font_size=18, bold=True, 
                              color=colors[i], align=PP_ALIGN.CENTER)

            self._add_text_box(slide, x + Inches(0.15), start_y + Inches(1.3), 
                              card_width - Inches(0.3), Inches(0.7),
                              example["description"], font_size=13, 
                              color=THEME.text_medium, align=PP_ALIGN.CENTER)

            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          x + Inches(0.3), start_y + Inches(2.1),
                                          card_width - Inches(0.6), Inches(0.5))
            badge.fill.solid()
            badge.fill.fore_color.rgb = THEME.primary_light
            badge.line.fill.background()

            self._add_text_box(slide, x + Inches(0.3), start_y + Inches(2.18),
                              card_width - Inches(0.6), Inches(0.4),
                              f"Benefit: {example['benefit']}", font_size=11, 
                              color=THEME.primary_dark, align=PP_ALIGN.CENTER)

    def create_comparison_slide(self, heading: str, pros: List[str], cons: List[str]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_background_shape(slide, THEME.bg_white)

        self._add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
                          heading, font_size=34, bold=True, color=THEME.primary_dark)

        vs_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          self.prs.slide_width / 2 - Inches(0.5),
                                          Inches(2.8), Inches(1), Inches(1))
        vs_circle.fill.solid()
        vs_circle.fill.fore_color.rgb = THEME.primary
        vs_circle.line.fill.background()

        vs_tf = vs_circle.text_frame
        vs_tf.word_wrap = False
        vs_p = vs_tf.paragraphs[0]
        vs_p.text = "VS"
        vs_p.font.size = Pt(28)
        vs_p.font.bold = True
        vs_p.font.color.rgb = THEME.text_white
        vs_p.alignment = PP_ALIGN.CENTER

        pros_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(0.8), Inches(1.3), Inches(4.5), Inches(0.6))
        pros_header.fill.solid()
        pros_header.fill.fore_color.rgb = THEME.accent_green
        pros_header.line.fill.background()

        self._add_text_box(slide, Inches(0.8), Inches(1.35), Inches(4.5), Inches(0.5),
                          "PROS", font_size=20, bold=True, 
                          color=THEME.text_white, align=PP_ALIGN.CENTER)

        y_pos = Inches(2.1)
        for pro in pros[:4]:
            check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y_pos + Inches(0.08),
                                          Inches(0.15), Inches(0.15))
            check.fill.solid()
            check.fill.fore_color.rgb = THEME.accent_green
            check.line.fill.background()

            self._add_text_box(slide, Inches(1.3), y_pos, Inches(3.8), Inches(0.6),
                              pro, font_size=15, color=THEME.text_medium)
            y_pos += Inches(0.7)

        cons_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(8), Inches(1.3), Inches(4.5), Inches(0.6))
        cons_header.fill.solid()
        cons_header.fill.fore_color.rgb = THEME.accent_red
        cons_header.line.fill.background()

        self._add_text_box(slide, Inches(8), Inches(1.35), Inches(4.5), Inches(0.5),
                          "CONS", font_size=20, bold=True, 
                          color=THEME.text_white, align=PP_ALIGN.CENTER)

        y_pos = Inches(2.1)
        for con in cons[:4]:
            x_mark = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), y_pos + Inches(0.08),
                                           Inches(0.15), Inches(0.15))
            x_mark.fill.solid()
            x_mark.fill.fore_color.rgb = THEME.accent_red
            x_mark.line.fill.background()

            self._add_text_box(slide, Inches(8.5), y_pos, Inches(3.8), Inches(0.6),
                              con, font_size=15, color=THEME.text_medium)
            y_pos += Inches(0.7)

    def create_summary_slide(self, heading: str, points: List[str]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_background_shape(slide, THEME.bg_light)

        top_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                           self.prs.slide_width, Inches(0.4))
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = THEME.primary
        top_accent.line.fill.background()

        self._add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
                          heading, font_size=36, bold=True, color=THEME.primary_dark)

        card_width = Inches(11.5)
        card_height = Inches(0.7)
        start_x = Inches(0.6)
        start_y = Inches(1.7)

        colors = [THEME.primary, THEME.accent_teal, THEME.accent_green, 
                 THEME.accent_orange, THEME.accent_purple]

        for i, point in enumerate(points[:5]):
            y = start_y + i * Inches(0.9)

            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          start_x, y, Inches(0.5), Inches(0.5))
            badge.fill.solid()
            badge.fill.fore_color.rgb = colors[i % len(colors)]
            badge.line.fill.background()

            self._add_text_box(slide, start_x, y + Inches(0.05), Inches(0.5), Inches(0.4),
                              str(i + 1), font_size=18, bold=True, 
                              color=THEME.text_white, align=PP_ALIGN.CENTER)

            point_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               start_x + Inches(0.6), y,
                                               card_width - Inches(0.6), card_height)
            point_card.fill.solid()
            point_card.fill.fore_color.rgb = THEME.bg_white
            point_card.line.color.rgb = THEME.bg_card

            self._add_text_box(slide, start_x + Inches(0.8), y + Inches(0.12),
                              card_width - Inches(1), Inches(0.5),
                              point, font_size=16, color=THEME.text_medium)

    def create_closing_slide(self, topic: str):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_background_shape(slide, THEME.primary)

        circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), 
                                        Inches(5), Inches(5))
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = RGBColor(0x4F, 0x46, 0xE5)
        circle1.line.fill.background()
        spTree = slide.shapes._spTree
        sp = circle1._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), 
                                        Inches(4), Inches(4))
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = RGBColor(0x4F, 0x46, 0xE5)
        circle2.line.fill.background()
        spTree = slide.shapes._spTree
        sp = circle2._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        self._add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
                          "Thank You!", font_size=60, bold=True, 
                          color=THEME.text_white, align=PP_ALIGN.CENTER)

        self._add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
                          f"Questions about {topic}?", font_size=28, 
                          color=THEME.primary_light, align=PP_ALIGN.CENTER)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     self.prs.slide_width / 2 - Inches(1.5), Inches(4.2),
                                     Inches(3), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = THEME.primary_light
        line.line.fill.background()

        self._add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
                          "Keep Learning - Keep Coding - Keep Growing", 
                          font_size=14, color=THEME.primary_light, align=PP_ALIGN.CENTER)

    def build(self, content: Dict, output_path: str):
        print(f"\nBuilding presentation: {content['title']}...")

        self.create_title_slide(content["title"], content["subtitle"])
        print("  [OK] Title slide")

        intro = content["introduction"]
        self.create_content_slide(intro["heading"], intro["points"], THEME.primary)
        print("  [OK] Introduction")

        why = content["why_it_matters"]
        self.create_content_slide(why["heading"], why["points"], THEME.accent_teal)
        print("  [OK] Why It Matters")

        core = content["core_concepts"]
        self.create_concept_slide(core["heading"], core["concepts"])
        print("  [OK] Core Concepts + Diagram")

        how = content["how_it_works"]
        self.create_flowchart_slide(how["heading"], how["steps"])
        print("  [OK] How It Works Flowchart")

        examples = content["real_world_examples"]
        self.create_examples_slide(examples["heading"], examples["examples"])
        print("  [OK] Real-World Examples")

        comp = content["comparison"]
        self.create_comparison_slide(comp["heading"], comp["pros"], comp["cons"])
        print("  [OK] Pros & Cons")

        summary = content["summary"]
        self.create_summary_slide(summary["heading"], summary["points"])
        print("  [OK] Summary")

        self.create_closing_slide(content["title"])
        print("  [OK] Closing slide")

        self.prs.save(output_path)
        print(f"\n[SUCCESS] Presentation saved!")
        print(f"File: {os.path.abspath(output_path)}")
        print(f"Slides: {len(self.prs.slides)}")


# =============================================================================
# MAIN
# =============================================================================

def generate_ppt(topic: str, output_path: str = None, use_api: bool = True) -> str:
    """
    Callable, server-friendly entry point.
    Takes a topic string and returns the filepath of the generated .pptx.
    Falls back to template mode automatically if no OPENAI_API_KEY is set.
    """
    topic = (topic or "").strip()
    if not topic:
        raise ValueError("No topic provided.")

    if use_api and not os.getenv("OPENAI_API_KEY"):
        use_api = False  # no key available on the server -> template mode

    if not output_path:
        safe_name = "".join(c if c.isalnum() or c in (" ", "-", "_") else "_"
                             for c in topic).strip().replace(" ", "_")
        output_path = f"{safe_name}.pptx"

    generator = ContentGenerator(use_api=use_api)
    content = generator.generate(topic)

    builder = PresentationBuilder()
    builder.build(content, output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Generate a modern CS presentation from a topic name.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python cs_pptx_generator.py "Data Structures"
  python cs_pptx_generator.py "Machine Learning" --output ml.pptx
  python cs_pptx_generator.py "Arrays" --no-api
        """
    )
    parser.add_argument("topic", help="Computer Science topic/subject name")
    parser.add_argument("--output", "-o", default=None, 
                       help="Output file path (default: <topic>.pptx)")
    parser.add_argument("--no-api", action="store_true",
                       help="Use template mode without AI (no API key needed)")

    args = parser.parse_args()

    if args.output:
        output_path = args.output
    else:
        safe_name = "".join(c if c.isalnum() or c in (" ", "-", "_") else "_" 
                           for c in args.topic).strip().replace(" ", "_")
        output_path = f"{safe_name}.pptx"

    if not args.no_api and not os.getenv("OPENAI_API_KEY"):
        print("No OPENAI_API_KEY found. Set it for AI-generated content.")
        print("  export OPENAI_API_KEY='your-key-here'")
        print("Or use --no-api for template mode.\n")
        response = input("Continue with template mode? (y/n): ").strip().lower()
        if response != "y":
            print("Exiting. Set your API key or use --no-api.")
            sys.exit(0)
        args.no_api = True

    print(f"\nTopic: {args.topic}")
    print("=" * 50)

    generator = ContentGenerator(use_api=not args.no_api)
    content = generator.generate(args.topic)

    builder = PresentationBuilder()
    builder.build(content, output_path)

    print("\nTip: Open the file in PowerPoint or Google Slides to present!")


if __name__ == "__main__":
    main()
